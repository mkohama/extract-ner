from pathlib import Path
from typing import List, Tuple, Any
from dataclasses import dataclass
from pptx import Presentation
from pptx.table import Table
from pptx.text.text import TextFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from langchain_core.documents import Document
from src.logger import log

# --- データクラス定義 ---


@dataclass
class ParsedShape:
    """パースされた図形の基本情報"""

    shape_type: str
    text: str
    left: int
    top: int
    width: int
    height: int
    row_num: int | None = None
    col_num: int | None = None


@dataclass
class ParsedSlide:
    """パースされたスライド"""

    slide_number: int
    shapes: List[ParsedShape]
    notes_text: str | None = None


@dataclass
class ParsedPresentation:
    """パースされたプレゼンテーション"""

    filename: str
    num_slides: int
    slide_width: int
    slide_height: int
    slides: List[ParsedSlide]


class PowerPointLoader:
    """
    PowerPointファイルローダー（Unstructured代替・軽量実装）

    特徴:
    - python-pptx のみ使用
    - テーブル抽出（CSV準拠のクォート処理と行列数の取得）
    - 画像やグループ図形も検出可能
    - スライド順・位置順にソート（スライド幅で正規化されたキーを使用）
    - LangChain Document対応
    """

    def __init__(
        self,
        file_path: str | Path,
        combine_slides: bool = False,
        notes_position: str = "bottom",
    ):
        """
        初期化

        Args:
            file_path: PowerPointファイルのパス
            combine_slides: Trueの場合、全スライドを1つのドキュメントにまとめる
            notes_position: スライドノートのコンテンツへの挿入位置 ("top"または"bottom"、デフォルトは"bottom")
        """
        self.file_path = Path(file_path)
        self.combine_slides = combine_slides
        self.notes_position = notes_position.lower()

        if not self.file_path.exists():
            raise FileNotFoundError(
                f"指定されたファイルが存在しません: {self.file_path}"
            )

        suffix = self.file_path.suffix.lower()
        if suffix == ".ppt":
            log(
                ".ppt 形式は python-pptx では正式対応していません。読み込みに失敗する可能性があります。"
            )
        elif suffix != ".pptx":
            raise ValueError(
                "サポートされていない形式です。現在は .pptx のみ対応しています。"
            )

    def load(self) -> List[Document]:
        """PowerPointファイルをLangChain Documentとして読み込む"""
        try:
            parsed = self.parse()
        except PackageNotFoundError:
            raise ValueError(
                f"PowerPointファイルが破損している可能性があります: {self.file_path}"
            )
        except Exception as e:
            log(f"PowerPointファイルの読み込み中にエラーが発生しました: {e}")
            raise RuntimeError(
                f"PowerPointファイルの読み込み中にエラーが発生しました: {e}"
            )

        # combine_slidesオプションに応じて出力
        if self.combine_slides:
            return self._to_combined_document(parsed)
        else:
            return self._to_documents(parsed)

    def parse(self) -> ParsedPresentation:
        """PowerPointファイルをパースして構造化データを返す"""
        presentation = Presentation(str(self.file_path))
        slides = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            # スライド幅をソートキー計算のために渡す
            shapes = self._parse_slide(slide, presentation.slide_width)

            # スライドノートを抽出
            notes_text = self._extract_slide_notes(slide)

            slides.append(
                ParsedSlide(
                    slide_number=slide_num, shapes=shapes, notes_text=notes_text
                )
            )

        # プレゼンテーションメタデータを付与して返す
        return ParsedPresentation(
            filename=str(self.file_path),
            num_slides=len(presentation.slides),
            slide_width=presentation.slide_width,
            slide_height=presentation.slide_height,
            slides=slides,
        )

    def _parse_slide(self, slide: Any, slide_width: int) -> List[ParsedShape]:
        """スライド内の図形を位置順にソートし、パースする"""
        shapes = []

        # 図形をソートし、抽出処理へ (グループ内の要素も正しい順序で抽出するため、ここでソートする)
        for shape in self._order_shapes(slide.shapes, slide_width):
            shapes.extend(self._extract_shape(shape, slide_width))
        return shapes

    def _extract_slide_notes(self, slide: Any) -> str | None:
        """スライドノートのテキストを抽出する"""
        if not slide.has_notes_slide:
            return None

        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame

        if not notes_text_frame:
            return None

        # _text_frame_to_textを利用してテキストを取得
        notes_text = self._text_frame_to_text(notes_text_frame)

        return notes_text if notes_text.strip() else None

    def _extract_shape(self, shape: Any, slide_width: int) -> List[ParsedShape]:
        """
        図形を解析し、必要に応じて再帰的に処理する。
        テーブルとテキストフレームの重複を防ぐためにif/elifを使用。
        """
        parsed_shapes = []

        try:
            # 1. テーブル図形
            if shape.has_table:
                # テーブルをCSV形式でパース
                text = self._table_to_text(shape.table)
                if text.strip():
                    parsed_shapes.append(
                        ParsedShape(
                            shape_type="table",
                            text=text,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                            row_num=len(shape.table.rows),
                            col_num=len(shape.table.columns),
                        )
                    )

            # 2. テキスト図形
            elif shape.has_text_frame:
                text = self._text_frame_to_text(shape.text_frame)
                if text.strip():
                    parsed_shapes.append(
                        ParsedShape(
                            shape_type="text",
                            text=text,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                        )
                    )

            # 3. グループ図形: 再帰的に処理
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # グループ内の図形を位置順にソートし、再帰的に抽出
                for sub_shape in self._order_shapes(shape.shapes, slide_width):
                    parsed_shapes.extend(self._extract_shape(sub_shape, slide_width))

            # 4. 画像図形: プレースホルダとして記録
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                parsed_shapes.append(
                    ParsedShape(
                        shape_type="image",
                        text="[Image]",
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                    )
                )

        except Exception as e:
            log(
                f"図形の解析中にエラーが発生しました（タイプ: {getattr(shape, 'shape_type', 'Unknown')}）: {e}"
            )

        return parsed_shapes

    def _to_documents(self, parsed: ParsedPresentation) -> List[Document]:
        """ParsedPresentationをスライドごとのDocumentに変換"""
        docs = []

        for slide in parsed.slides:
            content_parts = []
            notes_section = None

            if slide.notes_text:
                notes_section = f"[SLIDE NOTES]\n{slide.notes_text}"

            if slide.shapes:
                for shape in slide.shapes:
                    if shape.shape_type == "table":
                        # テーブルは行列数情報とCSVテキストとして挿入
                        content_parts.append(
                            f"[Table] ({shape.row_num} rows, {shape.col_num} cols)\n{shape.text}"
                        )
                    elif shape.shape_type == "image":
                        content_parts.append("[Image]")
                    else:
                        content_parts.append(shape.text)

            if notes_section:
                if self.notes_position == "top":
                    content_parts.insert(0, notes_section)
                else:
                    content_parts.append(notes_section)

            if not content_parts:
                continue

            content = "\n\n".join(content_parts)

            # メタデータの追加
            table_count = sum(1 for s in slide.shapes if s.shape_type == "table")
            image_count = sum(1 for s in slide.shapes if s.shape_type == "image")

            metadata = {
                "source": parsed.filename,
                "file_name": Path(parsed.filename).name,
                "file_type": "pptx",
                "slide_number": slide.slide_number,
                "total_slides": parsed.num_slides,
                "shapes_count": len(slide.shapes),
                "table_count": table_count,
                "image_count": image_count,
                "slide_width": parsed.slide_width,
                "slide_height": parsed.slide_height,
                "has_notes": bool(slide.notes_text),
                "notes_length": len(slide.notes_text) if slide.notes_text else 0,
            }

            docs.append(Document(page_content=content.strip(), metadata=metadata))

        return docs

    def _to_combined_document(self, parsed: ParsedPresentation) -> List[Document]:
        """全スライドを1つのDocumentにまとめる"""
        parts = []
        all_notes_text = []

        for slide in parsed.slides:
            if not slide.shapes and not slide.notes_text:
                continue

            # Combineモードの見栄え強化
            parts.append(f"\n=== Slide {slide.slide_number} ===")

            slide_parts = []
            notes_section = None

            if slide.notes_text:
                notes_section = f"[SLIDE NOTES]\n{slide.notes_text}"
                all_notes_text.append(slide.notes_text)

            # メインコンテンツの構築
            if slide.shapes:
                for shape in slide.shapes:
                    if shape.shape_type == "table":
                        slide_parts.append(
                            f"[Table] ({shape.row_num} rows, {shape.col_num} cols)\n{shape.text}"
                        )
                    elif shape.shape_type == "image":
                        slide_parts.append("[Image]")
                    else:
                        slide_parts.append(shape.text)

            # ノート位置に基づいてコンテンツに挿入
            if notes_section:
                if self.notes_position == "top":
                    parts.append(notes_section)
                    parts.extend(slide_parts)
                else:
                    parts.extend(slide_parts)
                    parts.append(notes_section)
            else:
                parts.extend(slide_parts)

            parts.append("")

        content = "\n".join(parts)

        # スライドノートのメタデータ化 (Combineモード用)
        has_any_notes = bool(all_notes_text)

        metadata = {
            "source": parsed.filename,
            "file_name": Path(parsed.filename).name,
            "file_type": "pptx",
            "format": "combined",
            "total_slides": parsed.num_slides,
            "has_notes": has_any_notes,
            "total_notes_length": sum(len(n) for n in all_notes_text),
        }

        return [Document(page_content=content.strip(), metadata=metadata)]

    def _text_frame_to_text(self, text_frame: TextFrame) -> str:
        """TextFrame内の文字列を結合（段落ごとに改行、空行を除去）"""

        #  空白段落の除去をより堅牢に
        paragraphs = [
            p.text.strip() for p in text_frame.paragraphs if p.text and p.text.strip()
        ]
        return "\n".join(paragraphs)

    def _table_to_text(self, table: Table) -> str:
        """テーブルをCSV準拠のテキストに変換（セル内クォート処理を含む）"""
        col_num = len(table.columns)
        rows = []
        row = []

        for i, cell in enumerate(table.iter_cells()):
            # 空セル対応
            if cell.text_frame is not None:
                # セル内のテキストを抽出・整形し、改行をスペースに置換
                text = (
                    self._text_frame_to_text(cell.text_frame).replace("\n", " ").strip()
                )
            else:
                text = ""

            # CSVクォート処理: カンマまたはダブルクォートが含まれている場合はクォートする
            if any(c in text for c in [",", '"']):
                # 埋め込みダブルクォートをエスケープ ('"' -> '""')
                text = '"' + text.replace('"', '""') + '"'

            row.append(text)

            # 行の終わり
            if (i + 1) % col_num == 0:
                rows.append(",".join(row))
                row = []

        return "\n".join(rows)

    def _order_shapes(self, shapes: Any, slide_width: int) -> List[Any]:
        """図形を自然な順序でソート（上→下、左→右）"""

        def sort_key(x: Any) -> Tuple[int, float]:
            # 丸め幅を8000に変更し、様々なスライドサイズでの安定性を向上
            row_group = round(x.top / 8000) * 8000
            # leftからの距離をスライド幅で正規化し、行内での相対位置を決定
            normalized_left = x.left / slide_width

            return (row_group, normalized_left)

        # topとleft属性を持つ図形のみをソート対象とする
        sortable_shapes = [
            s for s in shapes if hasattr(s, "top") and hasattr(s, "left")
        ]
        return sorted(sortable_shapes, key=sort_key)
